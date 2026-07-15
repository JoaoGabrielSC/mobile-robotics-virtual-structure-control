#!/usr/bin/env python3
"""Extract plain text from course PPTX slides into a single file.

Usage:
    uv run python scripts/extract_pptx.py
    uv run python scripts/extract_pptx.py --output docs/pptx_extract.txt
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_PATTERNS = [
    ROOT / "Part *.pptx",
    ROOT.parent / "Part 2-1 - Kinematic Models - Wheeled Mobile Robots.pptx",
]


def discover_pptx_files(root: Path) -> list[Path]:
    found: set[Path] = set()
    for path in sorted(root.glob("Part *.pptx")):
        found.add(path.resolve())
    parent = root.parent
    part_21 = parent / "Part 2-1 - Kinematic Models - Wheeled Mobile Robots.pptx"
    if part_21.exists():
        found.add(part_21.resolve())
    return sorted(found)


def shape_text(shape) -> list[str]:
    lines: list[str] = []
    if hasattr(shape, "text") and shape.text.strip():
        lines.append(shape.text.strip())
    if shape.has_table:
        for row in shape.table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                lines.append(row_text)
    return lines


def extract_presentation(path: Path) -> list[str]:
    lines: list[str] = []
    lines.append("\n" + "=" * 80)
    lines.append(f"FILE: {path.name}")
    lines.append("=" * 80)

    prs = Presentation(str(path))
    for slide_idx, slide in enumerate(prs.slides, 1):
        chunk: list[str] = []
        for shape in slide.shapes:
            chunk.extend(shape_text(shape))
        if chunk:
            lines.append(f"\n--- Slide {slide_idx} ---")
            lines.extend(chunk)
    return lines


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract text from course PPTX files")
    parser.add_argument(
        "--root",
        type=Path,
        default=ROOT,
        help="Directory containing Part *.pptx (default: project root)",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=ROOT / "docs" / "pptx_extract.txt",
        help="Output text file",
    )
    args = parser.parse_args()

    files = discover_pptx_files(args.root)
    if not files:
        raise SystemExit(f"No PPTX files found under {args.root}")

    output_lines: list[str] = [
        "PGEE5558 Mobile Robotics — PPTX text extraction",
        f"Source directory: {args.root}",
        f"Files: {len(files)}",
    ]
    for fp in files:
        output_lines.extend(extract_presentation(fp))

    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text("\n".join(output_lines), encoding="utf-8")
    print(f"Extracted {len(files)} presentations -> {args.output}")


if __name__ == "__main__":
    main()
